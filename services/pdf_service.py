import asyncio
import os
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple, BinaryIO
import img2pdf
from PIL import Image
import pypdf
from pypdf import PdfReader, PdfWriter, PdfMerger
import pdfplumber
from pdf2image import convert_from_path
import pikepdf
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from docx import Document
from docx.shared import Inches
import openpyxl
from pptx import Presentation
import io
import magic

from config import Config
from utils.helpers import generate_temp_filename, clean_temp_file

class PDFService:
    def __init__(self):
        self.temp_dir = Config.TEMP_DIR
        self.temp_dir.mkdir(parents=True, exist_ok=True)
    
    async def images_to_pdf(self, image_files: List[BinaryIO], output_path: Optional[str] = None) -> str:
        """Convert images to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        images = []
        for img_file in image_files:
            img = Image.open(img_file)
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            images.append(img)
        
        # Convert first image and append others
        if images:
            images[0].save(output_path, save_all=True, append_images=images[1:])
        
        return output_path
    
    async def merge_pdfs(self, pdf_files: List[str], output_path: Optional[str] = None) -> str:
        """Merge multiple PDF files"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        merger = PdfMerger()
        
        for pdf_file in pdf_files:
            merger.append(pdf_file)
        
        merger.write(output_path)
        merger.close()
        
        return output_path
    
    async def split_pdf(self, pdf_file: str, page_ranges: List[Tuple[int, int]], output_dir: Optional[str] = None) -> List[str]:
        """Split PDF by page ranges"""
        if not output_dir:
            output_dir = str(Config.TEMP_DIR)
        
        reader = PdfReader(pdf_file)
        output_files = []
        
        for i, (start, end) in enumerate(page_ranges):
            writer = PdfWriter()
            for page_num in range(start - 1, min(end, len(reader.pages))):
                writer.add_page(reader.pages[page_num])
            
            output_path = os.path.join(output_dir, f'split_{i+1}_{generate_temp_filename()}.pdf')
            with open(output_path, 'wb') as f:
                writer.write(f)
            output_files.append(output_path)
        
        return output_files
    
    async def compress_pdf(self, pdf_file: str, quality: str = 'medium', output_path: Optional[str] = None) -> str:
        """Compress PDF with quality levels"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        quality_settings = {
            'low': {'compress': True, 'quality': 30},
            'medium': {'compress': True, 'quality': 60},
            'high': {'compress': True, 'quality': 85}
        }
        
        settings = quality_settings.get(quality, quality_settings['medium'])
        
        # Use pikepdf for compression
        pdf = pikepdf.Pdf.open(pdf_file)
        pdf.save(output_path, 
                compress_streams=settings['compress'],
                object_stream_mode=pikepdf.ObjectStreamMode.generate)
        
        return output_path
    
    async def pdf_to_images(self, pdf_file: str, output_format: str = 'PNG', dpi: int = 200) -> List[str]:
        """Convert PDF pages to images"""
        images = convert_from_path(pdf_file, dpi=dpi)
        output_files = []
        
        for i, image in enumerate(images):
            output_path = os.path.join(Config.TEMP_DIR, f'page_{i+1}_{generate_temp_filename()}.{output_format.lower()}')
            image.save(output_path, format=output_format)
            output_files.append(output_path)
        
        return output_files
    
    async def pdf_to_word(self, pdf_file: str, output_path: Optional[str] = None) -> str:
        """Convert PDF to Word document"""
        if not output_path:
            output_path = generate_temp_filename('.docx')
        
        doc = Document()
        
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    doc.add_paragraph(text)
                
                # Extract tables
                tables = page.extract_tables()
                for table in tables:
                    if table:
                        doc.add_table(rows=len(table), cols=len(table[0]))
                        # Add table data
                
                doc.add_page_break()
        
        doc.save(output_path)
        return output_path
    
    async def word_to_pdf(self, word_file: str, output_path: Optional[str] = None) -> str:
        """Convert Word to PDF"""
        # This requires LibreOffice or similar
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        # Using python-docx and reportlab for conversion
        doc = Document(word_file)
        
        # Create PDF
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y_position = height - 50
        
        for paragraph in doc.paragraphs:
            if y_position < 50:
                c.showPage()
                y_position = height - 50
            
            c.drawString(50, y_position, paragraph.text[:100])  # Truncate long lines
            y_position -= 20
        
        c.save()
        return output_path
    
    async def rotate_pdf(self, pdf_file: str, rotation: int, pages: Optional[List[int]] = None, output_path: Optional[str] = None) -> str:
        """Rotate PDF pages"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            if pages is None or i + 1 in pages:
                page.rotate(rotation)
            writer.add_page(page)
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def add_watermark(self, pdf_file: str, watermark_text: str = None, watermark_image: str = None, 
                           output_path: Optional[str] = None) -> str:
        """Add text or image watermark to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        # Create watermark
        if watermark_text:
            watermark_pdf = self._create_text_watermark(watermark_text)
        elif watermark_image:
            watermark_pdf = self._create_image_watermark(watermark_image)
        else:
            return pdf_file
        
        watermark_page = PdfReader(watermark_pdf).pages[0]
        
        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        # Clean up
        os.unlink(watermark_pdf)
        
        return output_path
    
    def _create_text_watermark(self, text: str) -> str:
        """Create a text watermark PDF"""
        output_path = generate_temp_filename('.pdf')
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        c.saveState()
        c.translate(width/2, height/2)
        c.rotate(45)
        c.setFont("Helvetica", 40)
        c.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
        c.drawCentredString(0, 0, text)
        c.restoreState()
        c.save()
        
        return output_path
    
    def _create_image_watermark(self, image_path: str) -> str:
        """Create an image watermark PDF"""
        output_path = generate_temp_filename('.pdf')
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        c.saveState()
        c.translate(width/2, height/2)
        c.rotate(45)
        c.drawImage(image_path, -100, -100, width=200, height=200, mask='auto')
        c.restoreState()
        c.save()
        
        return output_path
    
    async def extract_images(self, pdf_file: str) -> List[str]:
        """Extract images from PDF"""
        output_files = []
        
        with pikepdf.Pdf.open(pdf_file) as pdf:
            for page_num, page in enumerate(pdf.pages):
                for image_key in page.images:
                    raw_image = page.images[image_key]
                    image = Image.open(io.BytesIO(raw_image.read()))
                    
                    output_path = os.path.join(
                        Config.TEMP_DIR,
                        f'image_page{page_num+1}_{generate_tempfilename()}.png'
                    )
                    image.save(output_path)
                    output_files.append(output_path)
        
        return output_files
    
    async def protect_pdf(self, pdf_file: str, password: str, output_path: Optional[str] = None) -> str:
        """Add password protection to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        writer.encrypt(password)
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def remove_password(self, pdf_file: str, password: str, output_path: Optional[str] = None) -> str:
        """Remove password from PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        
        if reader.is_encrypted:
            reader.decrypt(password)
        
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def rearrange_pages(self, pdf_file: str, page_order: List[int], output_path: Optional[str] = None) -> str:
        """Rearrange PDF pages"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for page_num in page_order:
            if 1 <= page_num <= len(reader.pages):
                writer.add_page(reader.pages[page_num - 1])
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def delete_pages(self, pdf_file: str, pages_to_delete: List[int], output_path: Optional[str] = None) -> str:
        """Delete selected pages from PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            if i + 1 not in pages_to_delete:
                writer.add_page(page)
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def extract_pages(self, pdf_file: str, pages_to_extract: List[int], output_path: Optional[str] = None) -> str:
        """Extract selected pages to new PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for page_num in pages_to_extract:
            if 1 <= page_num <= len(reader.pages):
                writer.add_page(reader.pages[page_num - 1])
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def repair_pdf(self, pdf_file: str, output_path: Optional[str] = None) -> str:
        """Attempt to repair damaged PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        try:
            # Try to open with pikepdf and save
            pdf = pikepdf.Pdf.open(pdf_file, allow_overwriting_input=True)
            pdf.save(output_path)
        except:
            # Fallback: try to read with pypdf
            reader = PdfReader(pdf_file, strict=False)
            writer = PdfWriter()
            
            for page in reader.pages:
                writer.add_page(page)
            
            with open(output_path, 'wb') as f:
                writer.write(f)
        
        return output_path
    
    async def add_page_numbers(self, pdf_file: str, output_path: Optional[str] = None) -> str:
        """Add page numbers to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            # Create overlay with page number
            packet = io.BytesIO()
            c = canvas.Canvas(packet, pagesize=(page.mediabox.width, page.mediabox.height))
            c.drawString(
                page.mediabox.width / 2,
                30,
                str(i + 1)
            )
            c.save()
            
            packet.seek(0)
            overlay = PdfReader(packet)
            
            page.merge_page(overlay.pages[0])
            writer.add_page(page)
        
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    
    async def excel_to_pdf(self, excel_file: str, output_path: Optional[str] = None) -> str:
        """Convert Excel to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        workbook = openpyxl.load_workbook(excel_file)
        sheet = workbook.active
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y_position = height - 50
        
        for row in sheet.iter_rows(values_only=True):
            if y_position < 50:
                c.showPage()
                y_position = height - 50
            
            row_text = '  |  '.join([str(cell) if cell else '' for cell in row])
            c.drawString(50, y_position, row_text[:150])  # Truncate
            y_position -= 20
        
        c.save()
        return output_path
    
    async def ppt_to_pdf(self, ppt_file: str, output_path: Optional[str] = None) -> str:
        """Convert PowerPoint to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        prs = Presentation(ppt_file)
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        for slide in prs.slides:
            y_position = height - 50
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50
                    
                    c.drawString(50, y_position, shape.text[:100])
                    y_position -= 20
            
            c.showPage()
        
        c.save()
        return output_path
    
    async def text_to_pdf(self, text_file: str, output_path: Optional[str] = None) -> str:
        """Convert text file to PDF"""
        if not output_path:
            output_path = generate_temp_filename('.pdf')
        
        with open(text_file, 'r', encoding='utf-8') as f:
            text = f.read()
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y_position = height - 50
        
        for line in text.split('\n'):
            if y_position < 50:
                c.showPage()
                y_position = height - 50
            
            c.drawString(50, y_position, line[:100])
            y_position -= 15
        
        c.save()
        return output_path
    
    async def pdf_to_text(self, pdf_file: str, output_path: Optional[str] = None) -> str:
        """Extract text from PDF"""
        if not output_path:
            output_path = generate_temp_filename('.txt')
        
        text = ''
        
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + '\n\n'
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        return output_path
    
    async def pdf_to_html(self, pdf_file: str, output_path: Optional[str] = None) -> str:
        """Convert PDF to HTML"""
        if not output_path:
            output_path = generate_temp_filename('.html')
        
        html_content = '<!DOCTYPE html>\n<html>\n<head>\n'
        html_content += '<meta charset="UTF-8">\n'
        html_content += '<title>PDF to HTML</title>\n'
        html_content += '<style>body { font-family: Arial; margin: 40px; }</style>\n'
        html_content += '</head>\n<body>\n'
        
        with pdfplumber.open(pdf_file) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    html_content += f'<p>{text.replace(chr(10), "<br>")}</p>\n'
        
        html_content += '</body>\n</html>'
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return output_path

# Global service instance
pdf_service = PDFService()
